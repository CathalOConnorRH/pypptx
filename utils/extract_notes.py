from pptx import Presentation

def extract_speaker_notes(pptx_file):
    prs = Presentation(pptx_file)
    notes = []

    for slide in prs.slides:
        try:
            # Check if the slide has a notes_slide
            if slide.has_notes_slide:
                note_slide = slide.notes_slide
                text = note_slide.notes_text_frame.text if note_slide.notes_text_frame else ""
            else:
                text = ""  # No notes for this slide
        except Exception as e:
            text = ""  # Catch any other unexpected issues and proceed
            print(f"Error extracting notes from slide: {e}")

        notes.append(text)

    return notes

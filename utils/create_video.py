from moviepy.audio.AudioClip import concatenate_audioclips
from moviepy.editor import ImageClip, AudioFileClip, concatenate_videoclips
import os

from moviepy.video.io.ImageSequenceClip import ImageSequenceClip
from pptx import Presentation
from natsort import natsorted

def create_video(pptx_file, output_file, slides_dir, audio_dir):
    prs = Presentation(pptx_file)
    clips = []

    for i, slide in enumerate(prs.slides):
        slide_image = f"{slides_dir}/slide_{i+1}.png"

        # Convert slide to image (export slide to image)
        slide.shapes._element.getparent().write(slide_image)  # Assumes slides are already exported to PNG

        audio_file = f"{audio_dir}/slide_{i+1}.mp3"
        image_clip = ImageClip(slide_image).set_duration(5)  # 5 sec per slide, can adjust duration
        audio_clip = AudioFileClip(audio_file)

        # Set the audio to the image
        image_clip = image_clip.set_audio(audio_clip)

        clips.append(image_clip)

    # Concatenate all the slides
    final_clip = concatenate_videoclips(clips, method="compose")
    final_clip.write_videofile(output_file, fps=24)

def create_video2(output_dir, slides_dir, audio_dir):
    # Gather all PNG images from the specified folder
    image_files = sorted([os.path.join(slides_dir, img) for img in os.listdir(slides_dir) if img.endswith('.png')])
    sorted_image_files = natsorted(image_files)

    if not image_files:
        print("No images found in the specified folder.")
        return

    # List to hold video clips
    video_clips = []

    # Create a video clip for each image and its corresponding audio
    for i, image_file in enumerate(sorted_image_files):
        print("")
        # Create a clip from the image
        image_clip = ImageSequenceClip([image_file], fps=24)  # Using a higher fps for smoother transitions

        # Load the corresponding audio file
        audio_file = os.path.join(audio_dir, f"slide_{i + 1}.mp3")  # Adjust the naming as per your convention
        if os.path.exists(audio_file):
            audio_clip = AudioFileClip(audio_file)
            # Set the audio to the image clip
            image_clip = image_clip.set_audio(audio_clip)
            # Set the duration of the clip to match the audio length
            image_clip = image_clip.set_duration(audio_clip.duration)
        else:
            print(f"No audio file found for slide {i + 1}, using default duration of 3 seconds.")
            image_clip = image_clip.set_duration(3)  # Default duration if audio not found

        video_clips.append(image_clip)

    # Concatenate all the video clips
    final_video = concatenate_videoclips(video_clips, method="compose")

    # Write the video file to the output path
    final_video.write_videofile(output_dir, codec='libx264', audio_codec='aac')
    print(f"Video created successfully at {output_dir}")
